"""InfraGenie agent-tier APIs: tenant summary (dashboard), reports, ITSM extensions.

Everything here consumes `tenant_ai_service.tenant_chat` or direct SP tools for accuracy.
Reports are generated with ReportLab and uploaded to Azure Blob storage under a folder
`infragenie-reports/<user_id>/`; download URLs are short-lived SAS links.
"""
from __future__ import annotations

import io
import json
import logging
import os
import uuid
from datetime import datetime, timezone, timedelta
from typing import Any, Dict, List, Optional

from fastapi import APIRouter, Depends, HTTPException
from pydantic import BaseModel

from tenant_ai_service import (
    TOOL_HANDLERS, _load_sp as load_sp, tenant_chat,
    _tool_list_resource_groups, _tool_list_resources, _tool_list_vms,
    _tool_get_costs, _tool_get_secure_score,
)

logger = logging.getLogger("agent_api")

_SUMMARY_CACHE: Dict[str, Dict[str, Any]] = {}  # user_id -> {"at": dt, "data": {...}}
_SUMMARY_TTL = 120  # seconds


class TicketCreate(BaseModel):
    title: str
    description: str = ""
    priority: str = "P3"
    category: str = "general"
    resource_ref: Optional[str] = None
    assignee_email: Optional[str] = None


class ReportGenerate(BaseModel):
    title: str
    prompt: str
    format: str = "pdf"  # pdf|csv|pptx

class ITSMChatIn(BaseModel):
    message: str

class ITSMChatOut(BaseModel):
    reply: str
    ticket: Optional[Dict[str, Any]] = None
    include_sections: List[str] = []  # optional pre-defined sections


class RemediateIn(BaseModel):
    policy_assignment_id: str
    note: str = ""


def create_agent_router(db, emit_event, get_current_user):
    router = APIRouter(prefix="/api")

    # -------- Tenant summary (cached) --------
    @router.get("/tenant/summary")
    async def tenant_summary(user: dict = Depends(get_current_user), fresh: bool = False):
        uid = user["id"]
        now = datetime.now(timezone.utc)
        cached = _SUMMARY_CACHE.get(uid)
        if cached and not fresh and (now - cached["at"]).total_seconds() < _SUMMARY_TTL:
            return {**cached["data"], "cached": True, "cached_at": cached["at"].isoformat()}
        user_full = await db.users.find_one({"id": uid}, {"_id": 0}) or user
        sp = await load_sp(db, user_full)
        if not sp:
            return {
                "configured": False, "resource_groups": 0, "resources": 0, "vms": 0,
                "monthly_cost": 0.0, "currency": "INR", "secure_score": 0, "cached": False,
            }
        import asyncio
        async def run(fn, **kw):
            return await asyncio.to_thread(fn, sp, **kw) if kw else await asyncio.to_thread(fn, sp)
        try:
            rg, resources, vms, costs, sec = await asyncio.gather(
                run(_tool_list_resource_groups),
                run(_tool_list_resources, top=500),
                run(_tool_list_vms),
                run(_tool_get_costs),
                run(_tool_get_secure_score),
                return_exceptions=True,
            )
        except Exception as e:  # noqa: BLE001
            raise HTTPException(502, f"Tenant summary failed: {e}")
        def _safe(v, path, default):
            if isinstance(v, Exception): return default
            try:
                cur = v
                for k in path: cur = cur[k]
                return cur
            except Exception:
                return default
        # resource type breakdown
        type_breakdown: Dict[str, int] = {}
        if not isinstance(resources, Exception):
            for r in resources.get("resources", []) or []:
                t = (r.get("type") or "unknown").split("/")[-1]
                type_breakdown[t] = type_breakdown.get(t, 0) + 1
        location_breakdown: Dict[str, int] = {}
        if not isinstance(resources, Exception):
            for r in resources.get("resources", []) or []:
                loc = r.get("location") or "unknown"
                location_breakdown[loc] = location_breakdown.get(loc, 0) + 1

        data = {
            "configured": True,
            "resource_groups": _safe(rg, ["count"], 0),
            "resource_groups_list": _safe(rg, ["resource_groups"], []),
            "resources": _safe(resources, ["count"], 0),
            "vms": _safe(vms, ["count"], 0),
            "vms_list": _safe(vms, ["vms"], []),
            "monthly_cost": round(_safe(costs, ["total"], 0.0), 2),
            "currency": _safe(costs, ["currency"], "INR"),
            "secure_score": _safe(sec, ["score"], 0),
            "type_breakdown": type_breakdown,
            "location_breakdown": location_breakdown,
            "generated_at": now.isoformat(),
        }
        _SUMMARY_CACHE[uid] = {"at": now, "data": data}
        return {**data, "cached": False}

    # -------- ITSM extensions --------
    @router.post("/itsm/tickets")
    async def create_ticket(payload: TicketCreate, user: dict = Depends(get_current_user)):
        from provisioning_service import new_ticket_number
        now = datetime.now(timezone.utc)
        tid = str(uuid.uuid4())
        tnum = new_ticket_number()
        ticket = {
            "id": tid, "ticket_number": tnum, "user_id": user["id"],
            "title": payload.title, "description": payload.description,
            "priority": payload.priority, "category": payload.category,
            "resource_ref": payload.resource_ref, "assignee_email": payload.assignee_email,
            "status": "open", "created_at": now, "updated_at": now,
            "audit": [{"action": "created", "by": user.get("email"), "at": now.isoformat(), "note": payload.description[:200]}],
            "comments": [], "source": "manual",
        }
        await db.tickets.insert_one(ticket)
        sn_synced = False
        try:
            from servicenow_service import create_incident as sn_create
            caller = payload.assignee_email or ""
            watch = "harshpardhi477@gmail.com"
            if caller and caller != watch:
                watch = f"{caller},{watch}"
            sn_result = await sn_create(
                db, user["id"],
                short_description=payload.title or "InfraGenie ticket",
                description=payload.description or "",
                caller_email=caller,
                watch_list=watch,
                severity=3 if payload.priority == "low" else 2 if payload.priority == "medium" else 1,
            )
            if sn_result:
                sn_synced = True
                await db.tickets.update_one({"id": tid}, {"$set": {"servicenow_sys_id": sn_result.get("sys_id"), "servicenow_number": sn_result.get("number")}})
        except Exception as e:
            logger.warning("ServiceNow forward failed: %s", e)
        await emit_event(user["id"], "ticket.created", f"Ticket {tnum} created", detail=payload.title[:80], level="info")
        return {"id": tid, "ticket_number": tnum, "servicenow_synced": sn_synced}

    @router.delete("/itsm/tickets/{ticket_id}")
    async def delete_ticket(ticket_id: str, destroy: bool = False, user: dict = Depends(get_current_user)):
        t = await db.tickets.find_one({"id": ticket_id, "user_id": user["id"]}, {"_id": 0})
        if not t:
            raise HTTPException(404, "Ticket not found")
        destroy_result = None
        if destroy and t.get("module_key"):
            # Reuse provisioning destroy
            from provisioning_service import stage_runtime_folder
            import terraform_runtime as tf
            user_full = await db.users.find_one({"id": user["id"]}, {"_id": 0}) or user
            sp = await load_sp(db, user_full)
            if sp:
                try:
                    tf_storage = await db.tf_storage_configs.find_one({"user_id": user["id"]}, {"_id": 0})
                    cwd = await stage_runtime_folder(user["id"], t["id"] + "-destroy", t["module_key"], t.get("collected_vars", {}))
                    await tf.tf_fmt(cwd)
                    init_code, _ = await tf.tf_init(cwd, sp, tf_storage)
                    if init_code == 0:
                        plan_code, _, _ = await tf.tf_plan(cwd, sp, destroy=True)
                        if plan_code == 0:
                            destroy_code, destroy_log = await tf.tf_destroy(cwd, sp)
                            destroy_result = {"status": "destroyed" if destroy_code == 0 else "failed", "log": destroy_log[-2000:] if destroy_log else ""}
                        else:
                            destroy_result = {"status": "failed", "log": "terraform plan (destroy) failed"}
                    else:
                        destroy_result = {"status": "failed", "log": "terraform init failed"}
                except Exception as e:  # noqa: BLE001
                    destroy_result = {"status": "failed", "error": str(e)[:300]}
        await db.tickets.delete_one({"id": ticket_id})
        await emit_event(user["id"], "ticket.deleted", f"Ticket {t.get('ticket_number')} deleted", detail="", level="warning")
        return {"deleted": True, "destroy": destroy_result}

    # -------- ITSM Chat Agent --------
    @router.post("/itsm/chat")
    async def itsm_chat(payload: ITSMChatIn, user: dict = Depends(get_current_user)):
        user_full = await db.users.find_one({"id": user["id"]}, {"_id": 0}) or user
        from provisioning_service import new_ticket_number

        system_prompt = (
            "You are an IT Service Management (ITSM) assistant integrated with ServiceNow. "
            "Your job is to help users create, track, and manage ServiceNow tickets through natural conversation.\n\n"
            "Available actions:\n"
            "1. CREATE TICKET — If the user asks to create a ticket/incident/request, collect: "
            "title/short_description (required), description, priority (P1-P4), category, assignee_email. "
            "After collecting enough details, call the CREATE_TICKET function.\n"
            "2. LIST TICKETS — If the user asks to see their tickets, call LIST_TICKETS.\n"
            "3. CHECK STATUS — If the user asks about a specific ticket's status, call GET_TICKET.\n\n"
            "Guidelines:\n"
            "- Be conversational and helpful. Ask clarifying questions if details are missing.\n"
            "- For ticket creation, ask for: title (required), description, priority (defaults to P3), category (defaults to general), assignee_email.\n"
            "- Always include harshpardhi477@gmail.com as a watcher on every ticket.\n"
            "- When a ticket is created in ServiceNow, confirm the ticket number to the user.\n"
            "- If ServiceNow is not configured, inform the user and create the ticket locally.\n"
            "- Keep responses concise.\n"
            "IMPORTANT: Your response must be valid JSON with keys: \"reply\" (your message to the user) and optionally \"action\" (\"create_ticket\", \"list_tickets\", \"get_ticket\", \"none\"). "
            "If action is create_ticket, also include a \"ticket\" object with the ticket details."
        )

        intent_prompt = (
            f"User message: {payload.message}\n\n"
            "Analyze the intent. If the user wants to create a ticket/incident/request, respond with a JSON object with keys: "
            "action (\"create_ticket\"), reply (your conversational response asking for missing details or confirming), "
            "and if all required fields (title) are present, include a ticket object with: "
            "title, description, priority, category, assignee_email.\n"
            "If they want to list tickets: action=\"list_tickets\", reply=\"...\"\n"
            "If they want to check a specific ticket: action=\"get_ticket\", reply=\"...\", ticket_number=\"...\"\n"
            "Otherwise: action=\"none\", reply=\"...\" (conversational only)\n\n"
            "Return ONLY valid JSON, no markdown."
        )

        try:
            result = await tenant_chat(
                db, user_full,
                [{"role": "system", "content": system_prompt}, {"role": "user", "content": intent_prompt}],
                max_tool_iterations=1,
            )
            raw = result["reply"]
            # Parse JSON from reply
            import re as _re
            m = _re.search(r'\{.*\}', raw, _re.DOTALL)
            parsed = json.loads(m.group(0)) if m else {"reply": raw, "action": "none"}
        except Exception:
            parsed = {"reply": "I understand you want help with ITSM. Could you clarify what you'd like to do — create a ticket, check status, or list your tickets?", "action": "none"}

        reply = parsed.get("reply", "")
        action = parsed.get("action", "none")
        ticket_data = parsed.get("ticket")

        # Execute action if needed
        created_ticket = None
        if action == "create_ticket" and ticket_data and ticket_data.get("title"):
            tid = str(uuid.uuid4())
            tnum = new_ticket_number()
            now = datetime.now(timezone.utc)
            title = ticket_data["title"]
            desc = ticket_data.get("description", "")
            priority = ticket_data.get("priority", "P3")
            category = ticket_data.get("category", "general")
            assignee = ticket_data.get("assignee_email", "")

            ticket_doc = {
                "id": tid, "ticket_number": tnum, "user_id": user["id"],
                "title": title, "description": desc,
                "priority": priority, "category": category,
                "assignee_email": assignee,
                "status": "open", "created_at": now, "updated_at": now,
                "audit": [{"action": "created", "by": user.get("email"), "at": now.isoformat(), "note": desc[:200]}],
                "comments": [], "source": "itsm_chat",
            }
            await db.tickets.insert_one(ticket_doc)

            sn_synced = False
            try:
                from servicenow_service import create_incident as sn_create
                watch = "harshpardhi477@gmail.com"
                if assignee and assignee != watch:
                    watch = f"{assignee},{watch}"
                sn_result = await sn_create(
                    db, user["id"],
                    short_description=title,
                    description=desc,
                    caller_email=assignee,
                    watch_list=watch,
                    severity=3 if priority == "P3" else 2 if priority in ("P2",) else 1,
                )
                if sn_result:
                    sn_synced = True
                    await db.tickets.update_one({"id": tid}, {"$set": {
                        "servicenow_sys_id": sn_result.get("sys_id"),
                        "servicenow_number": sn_result.get("number"),
                    }})
            except Exception as e:
                logger.warning("ServiceNow sync failed in ITSM chat: %s", e)

            created_ticket = {
                "id": tid, "ticket_number": tnum,
                "title": title, "status": "open",
                "servicenow_synced": sn_synced,
            }
            sn_ref = f" (ServiceNow: {sn_result.get('number')})" if sn_synced and sn_result else ""
            reply = f"Ticket **{tnum}** created for \"{title}\"{sn_ref}. I've added harshpardhi477@gmail.com as a watcher."

            await emit_event(user["id"], "ticket.created", f"Ticket {tnum} created via ITSM chat", detail=title[:80], level="info")

        elif action == "list_tickets":
            cursor = db.tickets.find({"user_id": user["id"]}).sort("created_at", -1).limit(20)
            tickets_list = []
            async for t in cursor:
                tickets_list.append(f"**{t.get('ticket_number')}** — {t.get('title', t.get('resource_name', '?'))} ({t.get('status')})")
            if tickets_list:
                reply = "Here are your recent tickets:\n" + "\n".join(tickets_list)
            else:
                reply = "You don't have any tickets yet. Would you like to create one?"

        return {"reply": reply, "ticket": created_ticket}

    # -------- Assessment agent --------
    @router.get("/assessments/catalog")
    async def assessment_catalog(user: dict = Depends(get_current_user)):
        return {"assessments": [
            {"key": "security", "label": "Security Posture", "desc": "Defender score, unassigned RBAC, exposed resources."},
            {"key": "cost", "label": "Cost & FinOps", "desc": "Idle resources, right-sizing, storage tiering."},
            {"key": "governance", "label": "Governance & Tags", "desc": "Tag coverage, naming, orphaned resources."},
            {"key": "reliability", "label": "Reliability", "desc": "SLA gaps, single-AZ resources, backup coverage."},
            {"key": "compliance", "label": "Compliance (GDPR/HIPAA/ISO)", "desc": "Policy assignment gaps by regulation."},
            {"key": "operations", "label": "Operational Excellence", "desc": "Monitoring, alerts, runbooks."},
        ]}

    @router.post("/assessments/{key}/run")
    async def run_assessment(key: str, user: dict = Depends(get_current_user)):
        user_full = await db.users.find_one({"id": user["id"]}, {"_id": 0}) or user
        hint_map = {
            "security": "Assess security posture; check Defender secure score, resources without encryption, public IPs, NSGs missing. Return findings JSON.",
            "cost": "Assess cost efficiency; list top spenders, idle resources, right-sizing opportunities. Return findings JSON.",
            "governance": "Assess governance; check tag coverage, naming inconsistencies, orphaned resources. Return findings JSON.",
            "reliability": "Assess reliability; check redundancy, backup coverage, single-AZ resources. Return findings JSON.",
            "compliance": "Assess compliance for GDPR/HIPAA/ISO27001; list policy gaps. Return findings JSON.",
            "operations": "Assess operational excellence; check monitoring, alert coverage, log analytics. Return findings JSON.",
        }
        hint = hint_map.get(key, f"Assess {key}")
        try:
            result = await tenant_chat(
                db, user_full,
                [{"role": "user", "content": f"Run a {key} assessment on my tenant. Use tools to fetch real data. Provide a scored assessment (0-100), summary, top findings, and recommended actions."}],
                hint=hint,
            )
            return {"key": key, "result": result["reply"], "tools": result.get("tool_traces", [])}
        except Exception as e:
            raise HTTPException(502, f"Assessment failed: {e}")

    # -------- Reports agent --------
    @router.get("/reports")
    async def list_reports(user: dict = Depends(get_current_user), limit: int = 30):
        cursor = db.reports.find({"user_id": user["id"]}, {"_id": 0}).sort("created_at", -1).limit(limit)
        items = []
        async for r in cursor:
            r.pop("body", None)
            items.append(r)
        return {"reports": items}

    @router.post("/reports/generate")
    async def generate_report(payload: ReportGenerate, user: dict = Depends(get_current_user)):
        user_full = await db.users.find_one({"id": user["id"]}, {"_id": 0}) or user
        now = datetime.now(timezone.utc)
        rid = str(uuid.uuid4())
        # 1) get AI content via tenant bridge (uses Azure Foundry if configured, else Azure OpenAI)
        foundry_cfg = await db.ai_configs.find_one({"user_id": user["id"], "provider": "foundry"}, {"_id": 0})
        provider_used = "azure_foundry" if foundry_cfg and foundry_cfg.get("api_key") else "azure_openai"
        prompt = f"Generate a professional executive report titled '{payload.title}'. User prompt: {payload.prompt}. Use tools to pull real tenant data. Structure: 1) Executive summary 2) Key metrics with numbers 3) Top findings 4) Recommendations. Be concise and data-driven."
        try:
            chat_res = await tenant_chat(db, user_full, [{"role": "user", "content": prompt}])
        except Exception as e:
            raise HTTPException(502, f"AI report generation failed: {e}")
        content = chat_res["reply"]

        # 2) render PDF
        pdf_bytes = _render_pdf(payload.title, user_full.get("email", ""), content, chat_res.get("tool_traces", []))

        # 3) upload to Azure blob (folder: infragenie-reports/<user_id>/)
        blob_url = None; download_link = None
        tf_cfg = await db.tf_storage_configs.find_one({"user_id": user["id"]}, {"_id": 0})
        if tf_cfg and tf_cfg.get("access_key"):
            try:
                blob_url, download_link = _upload_to_blob(tf_cfg, user["id"], rid, payload.title, pdf_bytes)
            except Exception as e:
                logger.warning("blob upload failed: %s", e)

        doc = {
            "id": rid, "user_id": user["id"], "title": payload.title, "format": payload.format,
            "provider": provider_used, "prompt": payload.prompt, "size_bytes": len(pdf_bytes),
            "created_at": now, "blob_url": blob_url, "download_link": download_link,
            "content_preview": content[:600], "content": content,
        }
        await db.reports.insert_one(doc)
        # also store binary for local fallback download
        await db.report_blobs.insert_one({"report_id": rid, "user_id": user["id"], "bytes": pdf_bytes})
        await emit_event(user["id"], "report.generated", f"Report '{payload.title}' ready", detail=f"{len(pdf_bytes)//1024} KB", level="success")
        # Return safe subset (exclude MongoDB _id and binary/large content field)
        return {
            "id": rid, "title": payload.title, "format": payload.format, "provider": provider_used,
            "size_bytes": len(pdf_bytes), "created_at": now.isoformat(),
            "blob_url": blob_url, "download_link": download_link,
            "content_preview": content[:600], "generated": True,
        }

    @router.get("/reports/{rid}/download")
    async def download_report(rid: str, user: dict = Depends(get_current_user)):
        r = await db.reports.find_one({"id": rid, "user_id": user["id"]}, {"_id": 0})
        if not r:
            raise HTTPException(404, "Report not found")
        if r.get("download_link"):
            return {"url": r["download_link"], "source": "azure_blob"}
        blob = await db.report_blobs.find_one({"report_id": rid, "user_id": user["id"]})
        if not blob:
            raise HTTPException(404, "Report file not found")
        import base64
        return {"url": None, "source": "local", "content_b64": base64.b64encode(blob["bytes"]).decode()}

    # -------- Policy remediation --------
    @router.post("/policy/remediate")
    async def policy_remediate(payload: RemediateIn, user: dict = Depends(get_current_user)):
        user_full = await db.users.find_one({"id": user["id"]}, {"_id": 0}) or user
        sp = await load_sp(db, user_full)
        if not sp:
            raise HTTPException(400, "Azure SP not configured")
        import asyncio
        def _remediate():
            from azure.mgmt.resource.policy import PolicyClient
            from tenant_ai_service import _sp_cred
            pc = PolicyClient(_sp_cred(sp), sp["subscription_id"])
            body = {"properties": {"policy_assignment_id": payload.policy_assignment_id}}
            name = f"remed-{uuid.uuid4().hex[:8]}"
            try:
                r = pc.remediations.create_or_update_at_subscription(remediation_name=name, parameters=body)
                return {"name": name, "id": r.id if hasattr(r, "id") else None}
            except Exception as e:
                raise
        try:
            result = await asyncio.to_thread(_remediate)
            await emit_event(user["id"], "policy.remediated", "Policy remediation triggered", detail=payload.policy_assignment_id[:80], level="success")
            return {"triggered": True, **result}
        except Exception as e:
            raise HTTPException(502, f"Remediation failed: {type(e).__name__}: {e}")

    @router.post("/reports/executive-deck")
    async def generate_executive_deck(user: dict = Depends(get_current_user)):
        """One-click 5-page executive board deck: Cost + Security + Reliability + Governance + Recommendations."""
        user_full = await db.users.find_one({"id": user["id"]}, {"_id": 0}) or user
        prompt = """Generate a 5-page executive board deck. Format each page starting with '## Page N: <title>' (exactly).

Page 1: Executive Summary — total resources, RGs, VMs, monthly cost (2 decimals with currency), secure score. 3-4 sentence overview.
Page 2: Cost Overview — pull top 10 resources by cost with get_costs, list total MTD spend, key trends, and 3 cost-optimization actions.
Page 3: Security Posture — pull get_secure_score and get_advisor_recommendations for Security. Show score, top 3 security risks, and remediation steps.
Page 4: Reliability & Governance — resource distribution by location/type, tag coverage assessment, single-AZ risks, backup gaps.
Page 5: Strategic Recommendations — 5 prioritized next actions with expected impact.

Use tools to pull real tenant data. Be data-driven, concise, board-ready. No fluff."""
        try:
            res = await tenant_chat(db, user_full, [{"role": "user", "content": prompt}], max_tool_iterations=6)
        except Exception as e:
            raise HTTPException(502, f"Deck generation failed: {e}")
        title = f"InfraGenie Executive Board Deck — {datetime.now(timezone.utc).strftime('%B %Y')}"
        pdf_bytes = _render_pdf(title, user_full.get("email", ""), res["reply"], res.get("tool_traces", []))
        rid = str(uuid.uuid4()); now = datetime.now(timezone.utc)
        blob_url = None; download_link = None
        tf_cfg = await db.tf_storage_configs.find_one({"user_id": user["id"]}, {"_id": 0})
        if tf_cfg and tf_cfg.get("access_key"):
            try: blob_url, download_link = _upload_to_blob(tf_cfg, user["id"], rid, title, pdf_bytes)
            except Exception as e: logger.warning("blob upload failed: %s", e)
        doc = {"id": rid, "user_id": user["id"], "title": title, "format": "pdf",
               "provider": "azure_openai", "prompt": "Executive Board Deck (auto)",
               "size_bytes": len(pdf_bytes), "created_at": now,
               "blob_url": blob_url, "download_link": download_link,
               "content_preview": res["reply"][:600], "content": res["reply"], "kind": "executive_deck"}
        await db.reports.insert_one(doc)
        await db.report_blobs.insert_one({"report_id": rid, "user_id": user["id"], "bytes": pdf_bytes})
        await emit_event(user["id"], "report.generated", f"Executive deck ready", detail=f"{len(pdf_bytes)//1024} KB", level="success")
        return {
            "id": rid, "title": title, "format": "pdf", "provider": "azure_openai",
            "size_bytes": len(pdf_bytes), "created_at": now.isoformat(),
            "blob_url": blob_url, "download_link": download_link,
            "content_preview": res["reply"][:600], "kind": "executive_deck",             "generated": True,
        }

    # -------- CSV/PPT Reports + Cache Invalidation --------
    @router.post("/reports/generate-csv")
    async def generate_csv_report(payload: ReportGenerate, user: dict = Depends(get_current_user)):
        """Generate a CSV report of infrastructure data."""
        user_full = await db.users.find_one({"id": user["id"]}, {"_id": 0}) or user
        now = datetime.now(timezone.utc)
        rid = str(uuid.uuid4())

        # Get AI content
        prompt = f"Generate a CSV-formatted report titled '{payload.title}'. User query: {payload.prompt}. Use tools to pull real data. Return only CSV data with a header row and data rows. No markdown, no explanation."
        try:
            chat_res = await tenant_chat(db, user_full, [{"role": "user", "content": prompt}])
        except Exception as e:
            raise HTTPException(502, f"AI report generation failed: {e}")
        csv_content = chat_res["reply"].strip()
        # Strip markdown fences if present
        if csv_content.startswith("```"):
            lines = csv_content.split("\n")
            lines = [l for l in lines if not l.startswith("```")]
            csv_content = "\n".join(lines).strip()

        csv_bytes = csv_content.encode("utf-8-sig")

        # Upload to blob
        blob_url = None; download_link = None
        tf_cfg = await db.tf_storage_configs.find_one({"user_id": user["id"]}, {"_id": 0})
        if tf_cfg and tf_cfg.get("access_key"):
            try:
                blob_url, download_link = _upload_to_blob_csv(tf_cfg, user["id"], rid, payload.title, csv_bytes)
            except Exception as e:
                logger.warning("csv blob upload failed: %s", e)

        doc = {
            "id": rid, "user_id": user["id"], "title": payload.title, "format": "csv",
            "provider": "azure_openai", "prompt": payload.prompt, "size_bytes": len(csv_bytes),
            "created_at": now, "blob_url": blob_url, "download_link": download_link,
            "content_preview": csv_content[:600], "content": csv_content,
        }
        await db.reports.insert_one(doc)
        await emit_event(user["id"], "report.generated", f"CSV Report '{payload.title}' ready", detail=f"{len(csv_bytes)//1024} KB", level="success")
        return {
            "id": rid, "title": payload.title, "format": "csv",
            "size_bytes": len(csv_bytes), "created_at": now.isoformat(),
            "blob_url": blob_url, "download_link": download_link,
            "content_preview": csv_content[:600], "generated": True,
        }

    @router.post("/reports/generate-pptx")
    async def generate_pptx_report(payload: ReportGenerate, user: dict = Depends(get_current_user)):
        """Generate a PowerPoint presentation from AI analysis."""
        user_full = await db.users.find_one({"id": user["id"]}, {"_id": 0}) or user
        now = datetime.now(timezone.utc)
        rid = str(uuid.uuid4())

        prompt = f"Generate a detailed analysis for a PowerPoint presentation titled '{payload.title}'. User query: {payload.prompt}. Use tools to pull real data. Structure the response as slides separated by '---SLIDE---'. Each slide should have a title on the first line after '---SLIDE---' prefixed with '## ', then content."
        try:
            chat_res = await tenant_chat(db, user_full, [{"role": "user", "content": prompt}])
        except Exception as e:
            raise HTTPException(502, f"AI report generation failed: {e}")
        content = chat_res["reply"]

        pptx_bytes = _render_pptx(payload.title, content)

        # Upload to blob
        blob_url = None; download_link = None
        tf_cfg = await db.tf_storage_configs.find_one({"user_id": user["id"]}, {"_id": 0})
        if tf_cfg and tf_cfg.get("access_key"):
            try:
                blob_url, download_link = _upload_to_blob_pptx(tf_cfg, user["id"], rid, payload.title, pptx_bytes)
            except Exception as e:
                logger.warning("pptx blob upload failed: %s", e)

        doc = {
            "id": rid, "user_id": user["id"], "title": payload.title, "format": "pptx",
            "provider": "azure_openai", "prompt": payload.prompt, "size_bytes": len(pptx_bytes),
            "created_at": now, "blob_url": blob_url, "download_link": download_link,
            "content_preview": content[:600], "content": content,
        }
        await db.reports.insert_one(doc)
        await emit_event(user["id"], "report.generated", f"PPT Report '{payload.title}' ready", detail=f"{len(pptx_bytes)//1024} KB", level="success")
        return {
            "id": rid, "title": payload.title, "format": "pptx",
            "size_bytes": len(pptx_bytes), "created_at": now.isoformat(),
            "blob_url": blob_url, "download_link": download_link,
            "content_preview": content[:600], "generated": True,
        }

    @router.post("/dashboard/invalidate-cache")
    async def invalidate_dashboard_cache(user: dict = Depends(get_current_user)):
        """Clear the cached tenant summary so it gets regenerated on next request."""
        _SUMMARY_CACHE.pop(user["id"], None)
        await emit_event(user["id"], "dashboard.cache.invalidated", "Dashboard cache cleared", detail="", level="info")
        return {"ok": True, "message": "Dashboard cache invalidated"}

    @router.post("/portal/logo")
    async def upload_logo(payload: dict, user: dict = Depends(get_current_user)):
        """Accept base64 logo data or URL and save it as portal_settings.logo_url."""
        logo_data = payload.get("logo_data") or payload.get("logo_url")
        if not logo_data:
            raise HTTPException(400, "logo_data or logo_url required")
        await db.portal_settings.update_one(
            {"user_id": user["id"]}, {"$set": {"logo_url": logo_data, "updated_at": datetime.now(timezone.utc)}}, upsert=True,
        )
        return {"ok": True, "logo_url": logo_data[:100]}

    # -------- Provisioning bundle mode --------
    @router.post("/provisioning/bundle/propose")
    async def propose_bundle(payload: dict, user: dict = Depends(get_current_user)):
        """Given a user intent, propose a bundle of Terraform modules (RG + VNet + Subnet + NSG + PIP + NIC + VM etc.)."""
        intent = payload.get("intent", "")
        if not intent.strip():
            raise HTTPException(400, "intent required")
        user_full = await db.users.find_one({"id": user["id"]}, {"_id": 0}) or user
        prompt = f"""User wants to deploy: {intent}

Available Terraform modules:
- resource-group, virtual-network, subnet, network-security-group, public-ip, network-interface (missing — synth if needed), virtual-machine-linux, virtual-machine-windows, storage-account, sql-server, sql-database, app-service, function-app, key-vault, load-balancer, managed-identity

Return STRICT JSON only:
{{
  "purpose": "<one-line summary>",
  "bundle_name": "<short name>",
  "modules": [
    {{"module_key": "resource-group", "reason": "container for all resources", "suggested_vars": {{"name":"rg-xxx","location":"centralindia"}}}},
    ...
  ],
  "estimated_monthly_cost_inr": <number>,
  "notes": "<what user should confirm before approving>"
}}

Rules:
- If they mention Linux VM, include: resource-group, virtual-network, subnet, network-security-group (with port 22 rule), public-ip, virtual-machine-linux
- If Windows VM: same but virtual-machine-windows with port 3389
- If web app: resource-group + app-service (+ optional key-vault, storage-account)
- If DB: resource-group + sql-server + sql-database
- Always suggest sensible defaults in suggested_vars."""
        try:
            res = await tenant_chat(db, user_full, [{"role": "user", "content": prompt}], max_tool_iterations=1)
        except Exception as e:
            raise HTTPException(502, f"Bundle proposal failed: {e}")
        # Try to parse JSON out of reply
        import re as _re
        text = res["reply"]
        m = _re.search(r"\{.*\}", text, _re.DOTALL)
        if m:
            try:
                parsed = json.loads(m.group(0))
                return {"intent": intent, "proposal": parsed, "raw": text[:400]}
            except Exception:
                pass
        return {"intent": intent, "proposal": None, "raw": text}

    @router.post("/provisioning/bundle/create")
    async def create_bundle_sessions(payload: dict, user: dict = Depends(get_current_user)):
        """Create one provisioning session per bundle module and return their IDs.

        The frontend can then walk through each session, or generate plans in a chain.
        """
        modules = payload.get("modules") or []
        if not modules:
            raise HTTPException(400, "modules[] required")
        now = datetime.now(timezone.utc)
        sessions = []
        for m in modules:
            key = m.get("module_key")
            if not key: continue
            sid = str(uuid.uuid4())
            await db.provisioning_sessions.insert_one({
                "id": sid, "user_id": user["id"], "workspace_id": user["id"],
                "module_key": key, "collected_vars": m.get("suggested_vars", {}),
                "missing_vars": [], "conversation": [
                    {"role": "assistant", "content": f"Bundle: {m.get('reason','')}", "at": now.isoformat()},
                ],
                "status": "ready", "created_at": now, "updated_at": now,
                "bundle_id": payload.get("bundle_name") or "bundle",
            })
            sessions.append({"session_id": sid, "module_key": key, "collected_vars": m.get("suggested_vars", {})})
        return {"sessions": sessions, "count": len(sessions)}

    # -------- Servicenow forwarding --------
    @router.get("/servicenow/health")
    async def servicenow_health(user: dict = Depends(get_current_user)):
        secret = await db.integration_secrets.find_one({"user_id": user["id"], "key": "servicenow"})
        conn = await db.integrations.find_one({"user_id": user["id"], "key": "servicenow"})
        if not (secret and conn and conn.get("connected")):
            return {"connected": False}
        import httpx
        url = conn.get("fields", {}).get("instance_url", "").rstrip("/")
        try:
            async with httpx.AsyncClient(auth=(conn["fields"].get("username", ""), secret.get("password", "")), timeout=10.0) as c:
                r = await c.get(f"{url}/api/now/table/incident?sysparm_limit=1")
                return {"connected": r.status_code == 200, "status_code": r.status_code}
        except Exception as e:
            return {"connected": False, "error": str(e)[:200]}

    # -------- Remote Script Execution via Azure VM Run Command --------
    @router.get("/scripts/catalog")
    async def script_catalog():
        """Return available troubleshooting scripts from the scripts/ directory."""
        from pathlib import Path
        scripts_root = Path(__file__).parent / "scripts"
        catalog = []
        for root, dirs, files in os.walk(str(scripts_root)):
            for f in files:
                if f.endswith((".sh", ".ps1")):
                    rel = os.path.relpath(os.path.join(root, f), str(scripts_root))
                    parts = rel.replace("\\", "/").split("/")
                    platform = parts[0] if len(parts) > 1 else "common"
                    catalog.append({
                        "path": rel.replace("\\", "/"),
                        "filename": f,
                        "platform": platform,
                        "name": f.replace(".sh", "").replace(".ps1", "").replace("_", " ").title(),
                    })
        return {"scripts": sorted(catalog, key=lambda x: x["name"])}

    @router.post("/scripts/estimate")
    async def estimate_script_impact(payload: dict, user: dict = Depends(get_current_user)):
        """Given a VM name and script, estimate impact and ask for approval."""
        vm_name = payload.get("vm_name", "")
        script_path = payload.get("script_path", "")
        resource_group = payload.get("resource_group", "")
        if not vm_name or not script_path:
            raise HTTPException(400, "vm_name and script_path required")
        user_full = await db.users.find_one({"id": user["id"]}, {"_id": 0}) or user
        sp = await load_sp(db, user_full)
        if not sp:
            raise HTTPException(400, "Azure SP not configured")
        try:
            vms_data = await asyncio.to_thread(_tool_list_vms, sp)
            vms = vms_data.get("vms", [])
            matching = [v for v in vms if vm_name.lower() in v.get("name", "").lower()]
            if not matching:
                raise HTTPException(404, f"VM '{vm_name}' not found in your Azure subscription")
            vm = matching[0]
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(502, f"Failed to look up VM: {e}")
        from pathlib import Path
        script_full = Path(__file__).parent / "scripts" / script_path
        if not script_full.exists():
            raise HTTPException(404, f"Script not found: {script_path}")
        preview = script_full.read_text(encoding="utf-8")[:2000]
        return {
            "vm": vm,
            "script": script_path,
            "preview": preview,
            "estimated_impact": "The script will execute via Azure VM Run Command. It may modify system configuration. Ensure you have a backup.",
            "requires_approval": True,
        }

    @router.post("/scripts/execute")
    async def execute_script(payload: dict, user: dict = Depends(get_current_user)):
        """Execute a script on a target VM via Azure VM Run Command."""
        vm_name = payload.get("vm_name", "")
        resource_group = payload.get("resource_group", "")
        script_path = payload.get("script_path", "")
        confirmed = payload.get("confirmed", False)
        if not confirmed:
            raise HTTPException(400, "Execution requires explicit confirmation (confirmed=true)")
        if not vm_name or not script_path:
            raise HTTPException(400, "vm_name, resource_group, and script_path required")

        from pathlib import Path
        script_full = Path(__file__).parent / "scripts" / script_path
        if not script_full.exists():
            raise HTTPException(404, f"Script not found: {script_path}")
        script_content = script_full.read_text(encoding="utf-8")

        user_full = await db.users.find_one({"id": user["id"]}, {"_id": 0}) or user
        sp = await load_sp(db, user_full)
        if not sp:
            raise HTTPException(400, "Azure SP not configured")

        if not resource_group:
            rgs_data = await asyncio.to_thread(_tool_list_resource_groups, sp)
            vms_data = await asyncio.to_thread(_tool_list_vms, sp)
            for vm in vms_data.get("vms", []):
                if vm_name.lower() in vm.get("name", "").lower():
                    vm_id = vm.get("id", "")
                    import re as _re
                    m = _re.search(r"/resourceGroups/([^/]+)", vm_id)
                    if m:
                        resource_group = m.group(1)
                    break

        if not resource_group:
            raise HTTPException(400, "Could not determine resource group — please provide it explicitly")

        try:
            def _run_command():
                from azure.mgmt.compute import ComputeManagementClient
                from tenant_ai_service import _sp_cred
                cred = _sp_cred(sp)
                cc = ComputeManagementClient(cred, sp["subscription_id"])
                poller = cc.virtual_machines.begin_run_command(
                    resource_group_name=resource_group,
                    vm_name=vm_name,
                    parameters={
                        "commandId": "RunShellScript" if not script_path.endswith(".ps1") else "RunPowerShellScript",
                        "script": [script_content],
                    }
                )
                result = poller.result(timeout=300)
                output = []
                if result.value:
                    for msg in result.value:
                        output.append({"message": msg.message or "", "code": msg.code or ""})
                return {"status": "completed", "output": output}
            result = await asyncio.to_thread(_run_command)
        except Exception as e:
            logger.exception("Run Command execution failed")
            result = {"status": "failed", "error": str(e)[:500]}

        from provisioning_service import new_ticket_number
        now = datetime.now(timezone.utc)
        tid = str(uuid.uuid4())
        tnum = new_ticket_number()
        ticket = {
            "id": tid, "ticket_number": tnum, "user_id": user["id"],
            "title": f"Script execution — {script_path} on {vm_name}",
            "description": f"Executed {script_path} on VM {vm_name} in {resource_group}",
            "module_key": "virtual-machine-linux" if not script_path.endswith(".ps1") else "virtual-machine-windows",
            "resource_name": vm_name,
            "operation": "troubleshoot",
            "status": result.get("status", "completed"),
            "logs": result.get("output", []),
            "outputs": {},
            "created_at": now,
            "updated_at": now,
            "audit": [{"action": "script_executed", "by": user.get("name", "User"), "at": now.isoformat(), "note": f"Script: {script_path}, VM: {vm_name}"}],
            "comments": [], "source": "troubleshoot",
        }
        await db.tickets.insert_one(ticket)
        await emit_event(user["id"], "provisioning.deployed" if result.get("status") == "completed" else "provisioning.error",
                         f"Script execution {'completed' if result.get('status') == 'completed' else 'failed'} — {script_path}",
                         detail=f"VM: {vm_name}", level="success" if result.get("status") == "completed" else "error", notify=True)
        return {
            "ticket_id": tid, "ticket_number": tnum,
            "vm": vm_name, "resource_group": resource_group,
            "script": script_path, "result": result,
        }

    # -------- AI Troubleshooting --------
    @router.post("/troubleshoot/analyze")
    async def troubleshoot_analyze(payload: dict, user: dict = Depends(get_current_user)):
        """AI-powered troubleshooting: user describes issue -> we query Azure Monitor -> recommend scripts."""
        issue = payload.get("issue", "").strip()
        vm_name = payload.get("vm_name", "")
        if not issue:
            raise HTTPException(400, "issue description required")

        user_full = await db.users.find_one({"id": user["id"]}, {"_id": 0}) or user
        sp = await load_sp(db, user_full)

        monitor_context = ""
        if sp and vm_name:
            try:
                def _get_metrics():
                    from azure.mgmt.monitor import MonitorManagementClient
                    from tenant_ai_service import _sp_cred
                    cred = _sp_cred(sp)
                    mm = MonitorManagementClient(cred, sp["subscription_id"])
                    now = datetime.now(timezone.utc)
                    start = now - timedelta(hours=6)
                    from azure.mgmt.compute import ComputeManagementClient
                    cc = ComputeManagementClient(cred, sp["subscription_id"])
                    for vm in cc.virtual_machines.list_all():
                        if vm_name.lower() in vm.name.lower():
                            vm_id = vm.id
                            import re as _re
                            m = _re.search(r"/subscriptions/[^/]+/resourceGroups/([^/]+)", vm_id)
                            rg = m.group(1) if m else None
                            if not rg:
                                return {"error": "Could not determine resource group"}
                            metrics_data = mm.metrics.list(
                                resource_uri=vm_id,
                                timespan=f"{start.isoformat()}/{now.isoformat()}",
                                interval="PT1H",
                                metricnames="Percentage CPU",
                                aggregation="Average,Maximum",
                            )
                            cpu_values = []
                            for t in metrics_data.value:
                                for series in t.timeseries:
                                    for point in series.data:
                                        if point.average is not None:
                                            cpu_values.append(round(point.average, 1))
                            return {"vm_name": vm.name, "resource_group": rg, "cpu_avg_6h": cpu_values[-6:] if len(cpu_values) >= 6 else cpu_values, "vm_size": vm.hardware_profile.vm_size if vm.hardware_profile else None}
                    return {"error": f"VM '{vm_name}' not found"}
                metrics = await asyncio.to_thread(_get_metrics)
                monitor_context = f"VM Metrics (last 6h): {json.dumps(metrics)[:1000]}"
            except Exception as e:
                monitor_context = f"Monitor query failed: {e}"

        from tenant_ai_service import tenant_chat
        prompt = f"""User reports issue: {issue}
VM: {vm_name}
Monitor context: {monitor_context}

Analyze the issue. Determine:
1. Root cause (most likely)
2. Recommended troubleshooting scripts from our catalog
3. Azure Monitor metrics to check
4. Estimated impact
5. Risk level (low/medium/high)

Return STRICT JSON:
{{
  "diagnosis": "<2-3 sentence root cause analysis>",
  "confidence": "high|medium|low",
  "recommended_scripts": [{{"path": "...", "reason": "..."}}],
  "monitor_metrics_to_check": ["..."],
  "estimated_impact": "<description>",
  "risk_level": "low|medium|high",
  "requires_approval": true
}}"""
        try:
            result = await tenant_chat(db, user_full, [{"role": "user", "content": prompt}], max_tool_iterations=3)
            import re as _re
            text = result["reply"]
            m = _re.search(r"\{.*\}", text, _re.DOTALL)
            if m:
                analysis = json.loads(m.group(0))
            else:
                analysis = {"diagnosis": text[:500], "confidence": "low", "recommended_scripts": [], "risk_level": "medium", "requires_approval": True}
        except Exception as e:
            analysis = {"diagnosis": f"Analysis failed: {e}", "confidence": "low", "recommended_scripts": [], "risk_level": "medium", "requires_approval": True}

        from pathlib import Path
        scripts_root = Path(__file__).parent / "scripts"
        available_scripts = []
        for root, dirs, files in os.walk(str(scripts_root)):
            for f in files:
                if f.endswith((".sh", ".ps1")):
                    rel = os.path.relpath(os.path.join(root, f), str(scripts_root)).replace("\\", "/")
                    available_scripts.append(rel)

        return {
            "issue": issue,
            "vm_name": vm_name,
            "analysis": analysis,
            "available_scripts": sorted(available_scripts),
            "monitor_context": monitor_context[:500] if monitor_context else None,
        }

    # -------- Monitoring / Alerting --------
    class AlertWebhookIn(BaseModel):
        message: str = ""
        data: Dict[str, Any] = {}

    class ApplyFixIn(BaseModel):
        ticket_number: str
        os_type: str = "linux"

    @router.post("/monitoring/webhook")
    async def monitoring_webhook(payload: AlertWebhookIn, user: dict = Depends(get_current_user)):
        from monitoring_service import handle_alert_webhook
        result = await handle_alert_webhook(db, user["id"], payload.model_dump())
        return result

    @router.post("/monitoring/apply-fix")
    async def monitoring_apply_fix(payload: ApplyFixIn, user: dict = Depends(get_current_user)):
        from monitoring_service import run_fix_script_on_vm
        ticket = await db.tickets.find_one({"ticket_number": payload.ticket_number, "user_id": user["id"]}, {"_id": 0})
        if not ticket:
            raise HTTPException(404, "Ticket not found")
        result = await run_fix_script_on_vm(
            db, user["id"], payload.ticket_number,
            ticket.get("subscription_id", ""),
            ticket.get("resource_group", ""),
            ticket.get("vm_name", ""),
            os_type=payload.os_type,
        )
        return result

    @router.get("/monitoring/alerts")
    async def monitoring_alerts(user: dict = Depends(get_current_user)):
        cursor = db.tickets.find(
            {"user_id": user["id"], "source": "monitoring"},
            {"_id": 0},
        ).sort("created_at", -1).limit(50)
        items = []
        async for t in cursor:
            items.append({
                "ticket_number": t.get("ticket_number", ""),
                "title": t.get("title", ""),
                "vm_name": t.get("vm_name", ""),
                "issue_type": t.get("issue_type", ""),
                "status": t.get("status", ""),
                "priority": t.get("priority", ""),
                "fix_applied": t.get("fix_applied", False),
                "created_at": t.get("created_at", ""),
                "servicenow_number": t.get("servicenow_number", ""),
                "metrics": t.get("metrics", {}),
            })
        return {"items": items, "total": len(items)}

    @router.post("/monitoring/alerts/{ticket_number}/dismiss")
    async def dismiss_alert(ticket_number: str, user: dict = Depends(get_current_user)):
        await db.tickets.update_one(
            {"ticket_number": ticket_number, "user_id": user["id"]},
            {"$set": {"status": "dismissed", "updated_at": datetime.now(timezone.utc)}},
        )
        from monitoring_service import ACTIVE_ALERTS
        for key in list(ACTIVE_ALERTS.keys()):
            if ACTIVE_ALERTS[key] == ticket_number:
                ACTIVE_ALERTS.pop(key, None)
        return {"dismissed": True}

    return router

def _render_pdf(title: str, user_email: str, body_md: str, tool_traces: list) -> bytes:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib.colors import HexColor
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
    import re

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4, leftMargin=2*cm, rightMargin=2*cm, topMargin=2*cm, bottomMargin=2*cm)
    styles = getSampleStyleSheet()
    accent = HexColor("#6366F1")

    title_style = ParagraphStyle("t", parent=styles["Heading1"], fontSize=22, textColor=accent, spaceAfter=6)
    meta_style = ParagraphStyle("m", parent=styles["Normal"], fontSize=9, textColor=HexColor("#64748B"), spaceAfter=12)
    body_style = ParagraphStyle("b", parent=styles["Normal"], fontSize=10.5, leading=15, textColor=HexColor("#0F172A"))
    h2 = ParagraphStyle("h2", parent=styles["Heading2"], fontSize=14, textColor=HexColor("#1E293B"), spaceBefore=14, spaceAfter=6)

    def _clean_inline(s: str) -> str:
        # Escape XML entities BEFORE inserting formatting tags
        s = s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        # Then convert markdown bold/italic to Paragraph tags
        s = re.sub(r"\*\*([^*]+?)\*\*", r"<b>\1</b>", s)
        s = re.sub(r"(?<![*])\*([^*]+?)\*(?![*])", r"<i>\1</i>", s)
        # Inline code -> monospace
        s = re.sub(r"`([^`]+?)`", r'<font name="Courier" size="9">\1</font>', s)
        return s

    story = []
    story.append(Paragraph("InfraGenie", ParagraphStyle("brand", parent=styles["Normal"], fontSize=11, textColor=accent, spaceAfter=2)))
    story.append(Paragraph(_clean_inline(title), title_style))
    story.append(Paragraph(f"Generated for {_clean_inline(user_email)} • {datetime.now(timezone.utc).strftime('%B %d, %Y %H:%M UTC')}", meta_style))

    for line in (body_md or "").split("\n"):
        s = line.rstrip()
        try:
            if not s:
                story.append(Spacer(1, 6))
            elif s.startswith("### "):
                story.append(Paragraph(_clean_inline(s[4:]), h2))
            elif s.startswith("## "):
                story.append(Paragraph(_clean_inline(s[3:]), h2))
            elif s.startswith("# "):
                story.append(Paragraph(_clean_inline(s[2:]), h2))
            elif s.startswith("- ") or s.startswith("* "):
                story.append(Paragraph("• " + _clean_inline(s[2:]), body_style))
            else:
                story.append(Paragraph(_clean_inline(s), body_style))
        except Exception as e:
            logger.warning("skipped bad line in PDF: %s", str(e)[:120])
            # Fallback: strip all special chars
            safe = "".join(c for c in s if c.isprintable() and c not in "<>&")
            try: story.append(Paragraph(safe, body_style))
            except Exception: pass

    if tool_traces:
        story.append(Spacer(1, 20))
        story.append(Paragraph("Data sources", h2))
        rows = [["Tool", "Args", "Bytes"]]
        for t in tool_traces:
            rows.append([str(t.get("tool", ""))[:30], json.dumps(t.get("args", {}))[:40], str(t.get("result_size", 0))])
        tbl = Table(rows, colWidths=[5*cm, 8*cm, 3*cm])
        tbl.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), accent),
            ("TEXTCOLOR", (0, 0), (-1, 0), HexColor("#FFFFFF")),
            ("FONTSIZE", (0, 0), (-1, -1), 8),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
            ("GRID", (0, 0), (-1, -1), 0.3, HexColor("#E2E8F0")),
        ]))
        story.append(tbl)

    story.append(Spacer(1, 30))
    story.append(Paragraph('<font color="#94A3B8" size="8">© InfraGenie — AI Cloud Operations for Azure</font>', body_style))
    doc.build(story)
    return buf.getvalue()


def _upload_to_blob(tf_cfg: dict, user_id: str, rid: str, title: str, data: bytes):
    from azure.storage.blob import BlobServiceClient, generate_blob_sas, BlobSasPermissions
    acct = tf_cfg["storage_account"]
    key = tf_cfg["access_key"]
    container = "infragenie-reports"
    conn_str = f"DefaultEndpointsProtocol=https;AccountName={acct};AccountKey={key};EndpointSuffix=core.windows.net"
    client = BlobServiceClient.from_connection_string(conn_str)
    # ensure container exists
    try:
        client.create_container(container)
    except Exception:
        pass
    safe_title = "".join(c if c.isalnum() or c in "._-" else "_" for c in title)[:60]
    blob_name = f"{user_id}/{datetime.now(timezone.utc).strftime('%Y%m%d')}-{rid[:8]}-{safe_title}.pdf"
    bc = client.get_blob_client(container=container, blob=blob_name)
    bc.upload_blob(data, overwrite=True, content_type="application/pdf")
    # generate 7-day SAS
    sas = generate_blob_sas(
        account_name=acct, container_name=container, blob_name=blob_name,
        account_key=key, permission=BlobSasPermissions(read=True),
        expiry=datetime.now(timezone.utc) + timedelta(days=7),
    )
    sas_url = f"https://{acct}.blob.core.windows.net/{container}/{blob_name}?{sas}"
    return bc.url, sas_url


def _render_pptx(title: str, body_md: str) -> bytes:
    """Render a PowerPoint presentation from markdown content (slides separated by ---SLIDE---)."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(99, 102, 241)  # accent indigo
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = f"InfraGenie \u2022 Generated {datetime.now(timezone.utc).strftime('%B %d, %Y')}"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(220, 220, 255)
    p2.alignment = PP_ALIGN.CENTER

    slides_data = body_md.split("---SLIDE---")
    for slide_text in slides_data:
        slide_text = slide_text.strip()
        if not slide_text:
            continue
        lines = slide_text.split("\n")
        slide_title = ""
        slide_content = []
        for line in lines:
            if line.startswith("## "):
                slide_title = line[3:].strip()
            elif line.startswith("# "):
                slide_title = line[2:].strip()
            else:
                slide_content.append(line)

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Title
        if slide_title:
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(1))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = slide_title
            p.font.size = Pt(32)
            p.font.color.rgb = RGBColor(99, 102, 241)
            p.font.bold = True
        # Content
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        y = 0
        for line in slide_content:
            line = line.strip()
            if not line:
                if y > 0:
                    p = tf.add_paragraph()
                    p.text = ""
                    p.font.size = Pt(14)
                y += 1
                continue
            if y == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            # Handle bold
            import re
            line_clean = re.sub(r"\*\*([^*]+?)\*\*", r"\1", line)
            p.text = line_clean
            p.font.size = Pt(16) if not line_clean.startswith("\u2022") else Pt(18)
            p.font.color.rgb = RGBColor(30, 41, 59)
            p.space_after = Pt(6)
            y += 1

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _upload_to_blob_csv(tf_cfg: dict, user_id: str, rid: str, title: str, data: bytes) -> tuple:
    """Upload CSV to Azure Blob, return (blob_url, sas_url)."""
    from azure.storage.blob import BlobServiceClient, generate_blob_sas, BlobSasPermissions
    acct = tf_cfg["storage_account"]
    key = tf_cfg["access_key"]
    container = "infragenie-reports"
    conn_str = f"DefaultEndpointsProtocol=https;AccountName={acct};AccountKey={key};EndpointSuffix=core.windows.net"
    client = BlobServiceClient.from_connection_string(conn_str)
    try:
        client.create_container(container)
    except Exception:
        pass
    safe_title = "".join(c if c.isalnum() or c in "._-" else "_" for c in title)[:60]
    blob_name = f"{user_id}/{datetime.now(timezone.utc).strftime('%Y%m%d')}-{rid[:8]}-{safe_title}.csv"
    bc = client.get_blob_client(container=container, blob=blob_name)
    bc.upload_blob(data, overwrite=True, content_type="text/csv")
    sas = generate_blob_sas(
        account_name=acct, container_name=container, blob_name=blob_name,
        account_key=key, permission=BlobSasPermissions(read=True),
        expiry=datetime.now(timezone.utc) + timedelta(days=7),
    )
    sas_url = f"https://{acct}.blob.core.windows.net/{container}/{blob_name}?{sas}"
    return bc.url, sas_url


def _upload_to_blob_pptx(tf_cfg: dict, user_id: str, rid: str, title: str, data: bytes) -> tuple:
    """Upload PPTX to Azure Blob, return (blob_url, sas_url)."""
    from azure.storage.blob import BlobServiceClient, generate_blob_sas, BlobSasPermissions
    acct = tf_cfg["storage_account"]
    key = tf_cfg["access_key"]
    container = "infragenie-reports"
    conn_str = f"DefaultEndpointsProtocol=https;AccountName={acct};AccountKey={key};EndpointSuffix=core.windows.net"
    client = BlobServiceClient.from_connection_string(conn_str)
    try:
        client.create_container(container)
    except Exception:
        pass
    safe_title = "".join(c if c.isalnum() or c in "._-" else "_" for c in title)[:60]
    blob_name = f"{user_id}/{datetime.now(timezone.utc).strftime('%Y%m%d')}-{rid[:8]}-{safe_title}.pptx"
    bc = client.get_blob_client(container=container, blob=blob_name)
    bc.upload_blob(data, overwrite=True, content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    sas = generate_blob_sas(
        account_name=acct, container_name=container, blob_name=blob_name,
        account_key=key, permission=BlobSasPermissions(read=True),
        expiry=datetime.now(timezone.utc) + timedelta(days=7),
    )
    sas_url = f"https://{acct}.blob.core.windows.net/{container}/{blob_name}?{sas}"
    return bc.url, sas_url
